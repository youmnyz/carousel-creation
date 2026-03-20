import copy
import os
from lxml import etree
from pptx import Presentation
from pptx.oxml.ns import qn
import config

# Tags mapped to carousel_data keys
TAG_MAP = {
    "{{TITLE}}":            "title",
    "{{HOOK}}":             "hook",
    "{{POINT_1_HEADLINE}}": "point_1_headline",
    "{{POINT_1_BODY}}":     "point_1_body",
    "{{POINT_2_HEADLINE}}": "point_2_headline",
    "{{POINT_2_BODY}}":     "point_2_body",
    "{{POINT_3_HEADLINE}}": "point_3_headline",
    "{{POINT_3_BODY}}":     "point_3_body",
    "{{POINT_4_HEADLINE}}": "point_4_headline",
    "{{POINT_4_BODY}}":     "point_4_body",
    "{{CONCLUSION}}":       "conclusion",
}

# These keys contain bullet-point text (lines separated by \n)
BULLET_KEYS = {
    "point_1_body", "point_2_body", "point_3_body", "point_4_body"
}


def _enable_autofit(shape):
    """Enable word wrap and shrink-to-fit on a shape's text frame."""
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    tf.word_wrap = True  # never break mid-word

    txBody = tf._txBody
    bodyPr = txBody.find(qn("a:bodyPr"))
    if bodyPr is None:
        return
    for child_tag in (qn("a:noAutofit"), qn("a:spAutoFit"), qn("a:normAutofit")):
        for child in bodyPr.findall(child_tag):
            bodyPr.remove(child)
    # normAutofit shrinks font size to fit — preserves layout, no word breaks
    etree.SubElement(bodyPr, qn("a:normAutofit"))


def _replace_run_with_bullets(run, value: str):
    """
    Replace run text with bullet lines separated by line breaks.
    Each line becomes a separate run; <a:br/> elements are inserted between them.
    Preserves the original run's formatting (rPr).
    """
    lines = [l for l in value.split("\n") if l.strip()]
    if not lines:
        run.text = ""
        return

    r_elem = run._r
    para_elem = r_elem.getparent()
    rPr = r_elem.find(qn("a:rPr"))

    # Set text of the original run to the first line
    run.text = lines[0]

    # Insert <a:br/> + new <a:r> for each subsequent line
    idx = list(para_elem).index(r_elem)
    for i, line in enumerate(lines[1:], start=1):
        # Line break element — copy rPr so spacing is consistent
        br = etree.Element(qn("a:br"))
        if rPr is not None:
            br.append(copy.deepcopy(rPr))
        para_elem.insert(idx + (i * 2 - 1), br)

        # New run with same formatting
        new_r = etree.Element(qn("a:r"))
        if rPr is not None:
            new_r.append(copy.deepcopy(rPr))
        new_t = etree.SubElement(new_r, qn("a:t"))
        new_t.text = line
        para_elem.insert(idx + (i * 2), new_r)


def _replace_tags_in_shape(shape, carousel_data: dict):
    """Find and replace all tags in a shape's text frame, then enable autofit."""
    if not shape.has_text_frame:
        return
    replaced = False
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            for tag, key in TAG_MAP.items():
                if tag in run.text:
                    value = str(carousel_data.get(key, ""))
                    if key in BULLET_KEYS:
                        # Strip the tag from run.text, then inject bullet lines
                        run.text = run.text.replace(tag, "")
                        _replace_run_with_bullets(run, value)
                    else:
                        run.text = run.text.replace(tag, value)
                    replaced = True
    if replaced:
        _enable_autofit(shape)


def _replace_tags_in_slide(slide, carousel_data: dict):
    for shape in slide.shapes:
        _replace_tags_in_shape(shape, carousel_data)
        # Handle grouped shapes
        if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
            for s in shape.shapes:
                _replace_tags_in_shape(s, carousel_data)


def build_carousel(post: dict, carousel_data: dict, output_path: str):
    template_path = config.TEMPLATE_FILE

    if not os.path.exists(template_path):
        raise FileNotFoundError(
            f"Template not found: {template_path}\n"
            f"Create a 6-slide template.pptx in the blog-to-carousel folder "
            f"with placeholder tags and set TEMPLATE_FILE in .env"
        )

    prs = Presentation(template_path)

    if len(prs.slides) != 6:
        raise ValueError(
            f"Template must have exactly 6 slides, found {len(prs.slides)}."
        )

    for slide in prs.slides:
        _replace_tags_in_slide(slide, carousel_data)

    os.makedirs(os.path.dirname(output_path) if os.path.dirname(output_path) else ".", exist_ok=True)
    prs.save(output_path)
    return output_path
